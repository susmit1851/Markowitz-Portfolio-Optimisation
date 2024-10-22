from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]

title.text = "Markowitz Portfolio Optimization"
subtitle.text = "Using MOSEK Fusion API"

# Slide 2: Overview
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_2.shapes.title
title.text = "Project Overview"
content = slide_2.shapes.placeholders[1].text = (
    "This project uses the Markowitz Portfolio Optimization model to maximize returns "
    "while imposing a limit on risk. The MOSEK Fusion API is used to solve this "
    "convex optimization problem with a no short-selling constraint."
)

# Slide 3: Mathematical Formulation
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_3.shapes.title
title.text = "Mathematical Formulation"

# LaTeX format equations explained in slide content
content = slide_3.shapes.placeholders[1]
content.text = (
    r"The Markowitz Portfolio Optimization problem is formulated as follows:\n\n"
    r"Maximize expected return:\n"
    r"$$ \text{maximize } \mu^T x $$\n\n"
    r"Subject to:\n"
    r"1. Budget constraint: $$ \sum_{i=1}^{n} x_i = 1 $$\n"
    r"2. Risk constraint: $$ x^T \Sigma x \leq \gamma $$\n"
    r"3. No short-selling constraint: $$ x_i \geq 0 $$ for all \(i\)\n\n"
)

# Slide 4: Optimization Problem
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_4.shapes.title
title.text = "Optimization Formulation"

content = slide_4.shapes.placeholders[1]
content.text = (
    r"The optimization problem can be written as:\n\n"
    r"$$\max_x \quad \mu^T x$$\n"
    r"Subject to:\n"
    r"1. Budget: $$\sum_{i=1}^{n} x_i = 1$$\n"
    r"2. Risk: $$x^T \Sigma x \leq \gamma$$\n"
    r"3. No short-selling: $$x_i \geq 0$$ for all \(i\)"
)

# Slide 5: Why Use a Quadratic Cone?
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_5.shapes.title
title.text = "Why Use a Quadratic Cone?"

content = slide_5.shapes.placeholders[1]
content.text = (
    r"The quadratic cone is used to impose the risk constraint, ensuring that the portfolio risk (variance) "
    r"remains below a certain threshold. This is done using a second-order cone constraint (Rotated QCone), which is "
    r"necessary because the risk constraint involves a quadratic term, $x^T \Sigma x$, which defines an ellipsoid in n-dimensional space.\n\n"
    r"The quadratic cone constraint ensures the model remains convex, allowing efficient solution of the optimization problem "
    r"using MOSEK."
)

# Slide 6: Python Code Overview
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_6.shapes.title
title.text = "Python Code Overview"
content = slide_6.shapes.placeholders[1]
content.text = (
    "1. Define mu (expected returns), sigma (covariance matrix), and gamma (risk tolerance).\n"
    "2. Use MOSEK API to set up variables and constraints.\n"
    "3. Define objective function to maximize return.\n"
    "4. Set risk constraint using a quadratic cone.\n"
    "5. Solve the model and retrieve results."
)

# Slide 7: Results and Conclusion
slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_7.shapes.title
title.text = "Results and Conclusion"
content = slide_7.shapes.placeholders[1]
content.text = (
    "The optimized portfolio allocates the fraction of holdings in each security "
    "such that the expected return is maximized while adhering to the budget and "
    "risk constraints. MOSEK provides efficient solutions for convex optimization."
)

# Save the presentation
prs.save("Markowitz_Optimization.pptx")

print("Presentation created successfully!")
